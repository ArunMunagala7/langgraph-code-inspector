#!/usr/bin/env python3
"""
Generate professional PowerPoint presentation from PPT_GUIDE.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
COLOR_PRIMARY = RGBColor(0, 102, 204)      # Blue
COLOR_ACCENT = RGBColor(0, 153, 76)       # Green
COLOR_DANGER = RGBColor(204, 0, 0)        # Red
COLOR_TEXT = RGBColor(51, 51, 51)         # Dark gray
COLOR_LIGHT = RGBColor(245, 245, 245)     # Light gray

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(60)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = subtitle
        subtitle_p.font.size = Pt(28)
        subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_PRIMARY
    title_shape.line.color.rgb = COLOR_PRIMARY
    
    # Title text
    title_frame = title_shape.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column comparison slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_PRIMARY
    title_shape.line.color.rgb = COLOR_PRIMARY
    
    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4.5), Inches(6))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    # Left title
    left_p = left_frame.paragraphs[0]
    left_p.text = f"❌ {left_title}"
    left_p.font.size = Pt(18)
    left_p.font.bold = True
    left_p.font.color.rgb = COLOR_DANGER
    
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1), Inches(4.2), Inches(6))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    # Right title
    right_p = right_frame.paragraphs[0]
    right_p.text = f"✅ {right_title}"
    right_p.font.size = Pt(18)
    right_p.font.bold = True
    right_p.font.color.rgb = COLOR_ACCENT
    
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(8)

# Slide 1: Title
add_title_slide(prs, 
    "Multi-Agent Code Understanding System",
    "From Generic Templates to Intelligent Code-Specific Analysis")

# Slide 2: Problem Statement
add_content_slide(prs, "The Problem",
    [
        "❌ All flowcharts were identical generic templates",
        "❌ No distinction between different algorithms",
        "❌ Conditions shown as generic 'IF/CONDITION CHECK'",
        "❌ Operations shown as generic 'Process/Compute'",
        "🔴 Root Cause: LLM prompt contained EXAMPLE template being copied"
    ])

# Slide 3: Solution Overview
add_content_slide(prs, "The Solution",
    [
        "✅ Intelligent Code-Aware Flowchart Generator v3",
        "✅ Direct Code Parsing (NO LLM for structure)",
        "✅ 3-Layer Fallback System",
        "✅ Actual Condition Details Extraction",
        "✅ Operation Context Visualization",
        "✅ Triple Sanitization Layer for Compatibility"
    ])

# Slide 4: Architecture Overview
add_content_slide(prs, "System Architecture",
    [
        "1. Direct Code Parsing",
        "   • Extract functions, loops, conditions, operations",
        "",
        "2. Intelligent Labeling (3-Layer Fallback)",
        "   • Layer 1: LLM Analysis | Layer 2: Regex Patterns | Layer 3: Default",
        "",
        "3. Build Mermaid Graph with Real Details",
        "",
        "4. Validate & Render with Auto-Correction"
    ])

# Slide 5: Before & After
add_two_column_slide(prs,
    "Before vs After Comparison",
    "BEFORE",
    [
        "Generic 'Process/Compute'",
        "Generic 'IF/CONDITION CHECK'",
        "All algorithms identical",
        "Low information value",
        "Cannot distinguish code"
    ],
    "AFTER",
    [
        "Specific labels like 'Find Two Sum'",
        "Actual conditions: 'arr[mid] == target'",
        "Unique per algorithm",
        "High information value",
        "Each algorithm clearly distinct"
    ])

# Slide 6: Test Results
add_content_slide(prs, "Test Results",
    [
        "✅ Two Sum: 'Find Two Sum' + 'Check- complement in seen'",
        "✅ Bubble Sort: 'Sort Algorithm' + 'Check- arr[j] > arr[j+1]'",
        "✅ Binary Search: 'Search Array' + Multiple conditions",
        "✅ LCS: 'Calculate LCS Table' + Nested loop conditions",
        "✅ Fibonacci: 'Recursive Compute' with recursion detection",
        "",
        "Test Coverage: 6 distinct algorithms | Success Rate: 100%"
    ])

# Slide 7: Key Improvements
add_content_slide(prs, "Key Improvements",
    [
        "📊 File Size: 60% reduction (53KB vs 149KB)",
        "🔷 Conditions: Show ACTUAL checks being performed",
        "📦 Operations: Display variable context & assignments",
        "🎯 Uniqueness: Each algorithm has distinct flowchart",
        "✅ Robustness: 3-layer fallback prevents failures",
        "⚡ Performance: <2s generation, <10s end-to-end"
    ])

# Slide 8: Three-Layer Fallback
add_content_slide(prs, "Intelligent Labeling: 3-Layer System",
    [
        "Layer 1: LLM Analysis (Primary)",
        "  → 'Find Two Sum', 'Sort Algorithm', 'Calculate LCS Table'",
        "",
        "Layer 2: Regex Patterns (Fallback 1)",
        "  → 'Swap/Update', 'Accumulate/Add', 'Recursive Compute'",
        "",
        "Layer 3: Safe Default (Fallback 2)",
        "  → 'Process/Compute' (always succeeds)"
    ])

# Slide 9: Condition Extraction
add_content_slide(prs, "Actual Condition Details",
    [
        "BEFORE: cond_3{\"IF/CONDITION CHECK\"}",
        "",
        "AFTER Examples:",
        "  • cond_2{\"Check- arr[mid] == target\"}",
        "  • cond_3{\"Check- arr[mid] < target\"}",
        "  • cond_2{\"Check- complement in seen\"}",
        "  • cond_3{\"Check- arr[j] > arr[j + 1]\"}"
    ])

# Slide 10: Operation Context
add_content_slide(prs, "Operation Context & Details",
    [
        "BEFORE: process_4[\"Process/Compute\"]",
        "",
        "AFTER Examples:",
        "  • process_4[\"Find Two Sum (seen = [])\"]",
        "  • process_4[\"Sort Algorithm (n = len(arr))\"]",
        "  • process_5[\"Search Array (left, right = 0...)\"]",
        "  • Includes actual variable assignments!"
    ])

# Slide 11: Technical Stack
add_content_slide(prs, "Technical Stack",
    [
        "Orchestration: LangGraph 1.0.4",
        "AI Engine: OpenAI GPT-4o-mini",
        "Visualization: Mermaid.js 11.12.0 + mmdc CLI",
        "Web UI: Gradio 6.0.2",
        "Language: Python 3.11+",
        "Database: None (lightweight JSON only)"
    ])

# Slide 12: Validation Pipeline
add_content_slide(prs, "Automatic Validation & Error Correction",
    [
        "1. Generate Mermaid syntax",
        "2. Validate with mmdc (15s timeout)",
        "",
        "If Invalid:",
        "  • LLM analyzes error message",
        "  • Attempts up to 3 automatic fixes",
        "  • Re-validates after each fix",
        "",
        "Success Rate: >95% | Fallback: matplotlib"
    ])

# Slide 13: Gradio Web Interface
add_content_slide(prs, "Gradio Web Interface",
    [
        "📝 Paste code or select from 10 test cases",
        "🔄 Real-time analysis with progress indicators",
        "📊 Tabbed results: Explanations | Analysis | Quality | Flowchart | Call Graph",
        "💾 Download diagrams (.png and .mmd files)",
        "🎨 Toggle between Mermaid and Matplotlib",
        "",
        "Access: python app.py → http://localhost:7860"
    ])

# Slide 14: Future Enhancements
add_content_slide(prs, "Future Enhancements",
    [
        "🌍 Multi-language support (JavaScript, Java, C++)",
        "🔍 Advanced condition parsing (complex expressions)",
        "⚡ Loop optimization detection",
        "📈 Complexity visualization on flowchart",
        "🎯 Interactive clickable nodes",
        "🔀 Side-by-side algorithm comparison",
        "📊 Performance profiling integration"
    ])

# Slide 15: Quick Start
add_content_slide(prs, "Quick Start Guide",
    [
        "1. git clone https://github.com/ArunMunagala7/langgraph-code-inspector.git",
        "2. cd langgraph-code-inspector",
        "3. python3 -m venv .venv",
        "4. source .venv/bin/activate",
        "5. pip install -r requirements.txt",
        "6. cp .env.example .env  # Add OpenAI API key",
        "7. python app.py  # Open http://localhost:7860"
    ])

# Slide 16: Conclusion
add_title_slide(prs,
    "Questions & Demo",
    "Production Ready • GitHub: ArunMunagala7/langgraph-code-inspector")

# Save presentation
output_path = "Multi_Agent_Code_Understanding_System.pptx"
prs.save(output_path)
print(f"✅ PowerPoint presentation created: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
